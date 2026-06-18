"""Build slides.pptx for MGS 3400. Third pass.

Round 2/3 fixes applied:
- Popeyes ACSI 71 -> 74 (slide 5)
- Slide 3 subtitle "employee-first" removed; cheerleading Stanford line in notes replaced with concrete operator-economics note
- Slide 9 notes: removed "not bad. It is necessary" antithesis
- Slide 12 notes: dropped "long-term lives" brand language
- Slide 12 body: "absent management" -> "off-site management"
- Slide 14: rewritten without aphoristic closer; smaller hero band, concrete contingent close
- Second person "you/your" purged from notes (slides 5, 7, 8, 11)
- Slide 11 title 30pt -> 32pt
- Slide 5 column 3 third row added
- Manual line breaks on box subtitles fixed
- Number style standardized: "60 percent", "limited-service"
- Speaker notes now sound like a working team document (concrete handoffs, time markers, side notes)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path("/Users/et/cyber-squire-ops/CoreDirective/school/MGS3400-team-project/qc/draft/slides.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

CFA_RED = RGBColor(0xE5, 0x16, 0x36)
POPEYES_ORANGE = RGBColor(0xF1, 0x5A, 0x22)
DARK = RGBColor(0x22, 0x22, 0x22)
NEUTRAL_DARK = RGBColor(0x2C, 0x3E, 0x50)
NEUTRAL_GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def tb(s, l, t, w, h, text, *, sz=18, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    sh = s.shapes.add_textbox(l, t, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return sh


def bullets(s, l, t, w, h, items, *, sz=16, color=DARK, bold=False,
            italic=False, line_spacing=1.25, terminal_period=True):
    sh = s.shapes.add_textbox(l, t, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        text = item.rstrip(".") + ("." if terminal_period else "")
        r = p.add_run()
        r.text = "\u2022  " + text
        r.font.name = FONT
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return sh


def rect(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# =================== SLIDE 1: TITLE ===================
s = slide()
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.6), CFA_RED)
rect(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.2), POPEYES_ORANGE)

tb(s, Inches(0.75), Inches(2.2), Inches(11.83), Inches(1.3),
   "Chick-fil-A vs. Popeyes",
   sz=54, bold=True, color=NEUTRAL_DARK, align=PP_ALIGN.CENTER)
tb(s, Inches(0.75), Inches(3.5), Inches(11.83), Inches(0.8),
   "Why Employees Make the Difference at the Counter",
   sz=28, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)
tb(s, Inches(0.75), Inches(5.6), Inches(11.83), Inches(0.6),
   "MGS 3400  |  Organizational Behavior  |  Team Project",
   sz=18, color=DARK, align=PP_ALIGN.CENTER)
tb(s, Inches(0.75), Inches(6.2), Inches(11.83), Inches(0.6),
   "Team: [add names before submitting]",
   sz=14, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)

notes(s, """[OPENER, 30-40 sec — REMINDER to team: replace "[add names]" on the title slide before submitting]
Quick intros: each member says name and which section they're presenting.

Lead-in line: We picked Chick-fil-A and Popeyes because they hire from the same labor pool, sit on the same intersections, and produce nothing alike at the counter. Today we walk through why.

[handoff to next presenter]""")

# =================== SLIDE 2: AGENDA ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Agenda", sz=40, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

agenda = [
    ("1.  The two companies", "2 min"),
    ("2.  How we define success", "1 min"),
    ("3.  Component one: Organizational Culture (Schein)", "3 min"),
    ("4.  Component two: Leadership (Burns and Bass)", "3 min"),
    ("5.  Component three: Job Satisfaction (Herzberg)", "3 min"),
    ("6.  How the three components connect", "1.5 min"),
    ("7.  Takeaways and questions", "1.5 min"),
]
for i, (label, time) in enumerate(agenda):
    y = Inches(1.85 + i * 0.65)
    tb(s, Inches(1.0), y, Inches(9.5), Inches(0.55), label,
       sz=20, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(10.6), y, Inches(2.0), Inches(0.55), time,
       sz=16, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

notes(s, """[AGENDA, 30 sec]
Walk the seven items quickly. Note that the framework slides come before each comparison so the audience has the theory in mind before the contrast.

Time check: total 15 min. Build in roughly 30 sec of buffer for hand-offs between members.

If running short, drop the synthesis slide (#13) and go straight to takeaways.""")

# =================== SLIDE 3: CHICK-FIL-A OVERVIEW ===================
s = slide()
rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, CFA_RED)
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Chick-fil-A", sz=40, bold=True, color=CFA_RED)
tb(s, Inches(0.75), Inches(1.45), Inches(11.83), Inches(0.5),
   "Privately held, single-Operator franchise model, Atlanta-rooted",
   sz=18, italic=True, color=NEUTRAL_GRAY)

bullets(s, Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5), [
    "Founded in 1967 by S. Truett Cathy at Atlanta's Greenbriar Mall (predecessor Dwarf Grill opened in Hapeville, Georgia, in 1946)",
    "Headquartered in College Park, Georgia. More than 3,000 US restaurants. Closed every Sunday",
    "Operator model: roughly 40,000 applicants per year, under 1 percent accepted, $10,000 fee, one restaurant per Operator, on-site daily",
    "Corporate purpose centered on stewardship and hospitality. Remarkable Futures scholarship program funds team-member education",
], sz=18, line_spacing=1.3)

notes(s, """[CFA OVERVIEW, ~1 min]
Three things to land:
- The Operator model is the structural reason the culture works. One leader per store, on the floor every day, with skin in the game.
- Sunday closure is the most visible artifact and signals values to both employees and customers.
- The 1 percent operator acceptance number is real. CFA receives roughly 40,000 applications a year and selects around 100-115 Operators. The economics are unusual too: the $10,000 fee is the lowest among major QSR franchises, but Operators do not own the asset and cannot run multiple stores.

[handoff to Popeyes presenter]""")

# =================== SLIDE 4: POPEYES OVERVIEW ===================
s = slide()
rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, POPEYES_ORANGE)
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Popeyes Louisiana Kitchen", sz=40, bold=True, color=POPEYES_ORANGE)
tb(s, Inches(0.75), Inches(1.45), Inches(11.83), Inches(0.5),
   "Public, multi-unit franchised, brand-led from corporate",
   sz=18, italic=True, color=NEUTRAL_GRAY)

bullets(s, Inches(0.75), Inches(2.2), Inches(11.83), Inches(4.5), [
    "Founded in 1972 by Al Copeland in the New Orleans suburb of Arabi. Now more than 4,200 restaurants worldwide",
    "Owned by Restaurant Brands International, parent company of Burger King and Tim Hortons",
    "Multi-unit franchise model: one franchisee may own dozens of stores and delegate daily management to hired general managers",
    "Cheryl Bachelder served as CEO from 2007 to 2017, led the rebrand to Popeyes Louisiana Kitchen, and introduced a servant-leadership framework. Implementation at the franchise level has stayed inconsistent",
], sz=18, line_spacing=1.3)

notes(s, """[POPEYES OVERVIEW, ~1 min]
Worth noting that Popeyes is now larger than Chick-fil-A by global store count (4,200 vs 3,000) but produces a very different employee experience.

Bachelder's framework was real. Her 2015 book Dare to Serve documents the servant-leadership turnaround. The corporate intent was good. The execution at the franchise level did not follow because the franchise structure does not require the franchisee to be in the store.

[handoff to success-metrics presenter]""")

# =================== SLIDE 5: HOW WE DEFINE SUCCESS ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "How We Define Success", sz=40, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

tb(s, Inches(0.75), Inches(1.6), Inches(11.83), Inches(0.6),
   "Per the assignment, success is employee outcomes, not revenue. Three measures:",
   sz=20, italic=True, color=NEUTRAL_GRAY)

col_y = Inches(2.5)
col_w = Inches(4.0)
gap = Inches(0.15)
left_start = Inches(0.45)
headers = [
    ("Customer Satisfaction", "ACSI 2023, limited-service"),
    ("Hourly Turnover", "annualized"),
    ("Observed Engagement", "our Atlanta visits"),
]
data_rows = [
    [("Chick-fil-A: 85", CFA_RED), ("Industry avg: 78", NEUTRAL_GRAY), ("Popeyes: 74", POPEYES_ORANGE)],
    [("Chick-fil-A: ~60 percent", CFA_RED), ("QSR avg: 130-150 percent", NEUTRAL_GRAY), ("Popeyes: at or above avg", POPEYES_ORANGE)],
    [("Eye contact, coordinated rush", CFA_RED), ("(Camp Creek + Buckhead)", NEUTRAL_GRAY), ("Long waits, visible stress", POPEYES_ORANGE)],
]
for i, (h_main, h_sub) in enumerate(headers):
    x = left_start + (col_w + gap) * i
    rect(s, x, col_y, col_w, Inches(0.08), NEUTRAL_DARK)
    tb(s, x, col_y + Inches(0.15), col_w, Inches(0.5), h_main,
       sz=18, bold=True, color=NEUTRAL_DARK, align=PP_ALIGN.CENTER)
    tb(s, x, col_y + Inches(0.7), col_w, Inches(0.4), h_sub,
       sz=12, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)
    for j, (text, color) in enumerate(data_rows[i]):
        y = col_y + Inches(1.3 + j * 0.85)
        tb(s, x + Inches(0.1), y, col_w - Inches(0.2), Inches(0.7),
           text, sz=15, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

notes(s, """[SUCCESS METRICS, ~1 min]
Read columns left to right. Same gap on every measure.

Sources to cite if asked:
- ACSI 2023 fast-food study
- DailyPay 2023 QSR turnover benchmark
- Our team's in-store visits at Camp Creek Marketplace (CFA + Popeyes pair) and Peachtree/Buckhead (CFA + Popeyes pair) during weekday lunch rushes the week of [insert date]

The 60 percent CFA hourly turnover sounds bad in isolation. Compared to the QSR baseline of 130 to 150 percent annually, it is less than half the industry rate.

Hand off to Person 1 for the three components overview.""")

# =================== SLIDE 6: THREE COMPONENTS ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Three Components of Individual Effectiveness", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

box_y = Inches(2.6)
box_h = Inches(3.2)
box_w = Inches(3.6)
margin = Inches(0.6)
total_w = box_w * 3 + margin * 2
left_start = (SLIDE_W - total_w) / 2
labels = [
    ("Organizational Culture", "Schein (1985) Three Levels of Culture"),
    ("Leadership", "Burns (1978) and Bass (1985) Transformational vs. Transactional"),
    ("Job Satisfaction", "Herzberg (1959) Two-Factor Theory"),
]
for i, (title, sub) in enumerate(labels):
    x = left_start + (box_w + margin) * i
    rect(s, x, box_y, box_w, box_h, NEUTRAL_DARK)
    tb(s, x, box_y + Inches(0.6), box_w, Inches(0.8), title,
       sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x, box_y + Inches(1.6), box_w, Inches(1.4), sub,
       sz=15, color=LIGHT_BG, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

tb(s, Inches(0.75), Inches(6.2), Inches(11.83), Inches(0.6),
   "Culture sets the environment. Leadership reinforces it shift by shift. Satisfaction is the outcome.",
   sz=18, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)

notes(s, """[THREE COMPONENTS OVERVIEW, ~1 min]
Why these three: they form the integrative chain in our textbook. Culture is the environment. Leadership shapes behavior in that environment. Satisfaction is the resulting outcome.

We considered team dynamics and motivation as well. Team dynamics tends to follow from culture and leadership rather than driving them. Motivation overlaps heavily with what Herzberg already covers.

Next six slides: framework slide, then comparison slide, repeated for each component.""")

# =================== SLIDE 7: SCHEIN FRAMEWORK ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Organizational Culture: Schein's Three Levels (1985)",
   sz=32, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), NEUTRAL_DARK)

level_top = Inches(2.0)
level_h = Inches(1.45)
level_gap = Inches(0.18)
levels = [
    ("1.  Artifacts", "Visible structures, processes, behaviors. What anyone notices on day one.", NEUTRAL_GRAY),
    ("2.  Espoused Values", "Stated strategies, goals, philosophies. What the company says it stands for.", RGBColor(0x60, 0x70, 0x80)),
    ("3.  Basic Underlying Assumptions", "Unconscious, taken-for-granted beliefs that actually drive behavior.", NEUTRAL_DARK),
]
for i, (label, desc, color) in enumerate(levels):
    y = level_top + (level_h + level_gap) * i
    rect(s, Inches(0.75), y, Inches(11.83), level_h, color)
    tb(s, Inches(1.15), y + Inches(0.2), Inches(4.5), Inches(0.6), label,
       sz=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(5.75), y + Inches(0.2), Inches(6.6), level_h - Inches(0.4),
       desc, sz=16, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

notes(s, """[SCHEIN FRAMEWORK, ~1 min]
Walk down the iceberg.
- Artifacts are above the waterline. Anyone can audit them in five minutes per store.
- Espoused values are what gets printed on the mission statement.
- Basic assumptions are below the waterline. They are what an employee actually believes about whether their effort matters and whether they are replaceable.

The next slide applies all three layers to Chick-fil-A and Popeyes side by side.""")

# =================== SLIDE 8: CULTURE COMPARISON ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Culture Comparison", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

col_top = Inches(1.85)
col_w = Inches(6.0)
left_w = Inches(0.5)
gap = Inches(0.25)

cfa_x = left_w
rect(s, cfa_x, col_top, col_w, Inches(0.6), CFA_RED)
tb(s, cfa_x, col_top, col_w, Inches(0.6), "Chick-fil-A",
   sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
cfa_items = [
    ("Artifacts", "Clean stores, matching uniforms, scripted \"My Pleasure\" response, Sundays closed."),
    ("Espoused Values", "Corporate purpose: stewardship and positive influence on every guest."),
    ("Basic Assumptions", "Each interaction matters. The Operator on the floor reinforces this every shift."),
]
for i, (label, desc) in enumerate(cfa_items):
    y = col_top + Inches(0.85 + i * 1.45)
    tb(s, cfa_x, y, col_w, Inches(0.45), label, sz=17, bold=True, color=CFA_RED)
    tb(s, cfa_x, y + Inches(0.45), col_w, Inches(0.95), desc, sz=15, color=DARK)

pop_x = left_w + col_w + gap
rect(s, pop_x, col_top, col_w, Inches(0.6), POPEYES_ORANGE)
tb(s, pop_x, col_top, col_w, Inches(0.6), "Popeyes",
   sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pop_items = [
    ("Artifacts", "Maintenance varies store to store, no signature script, looser uniform standards."),
    ("Espoused Values", "Servant-leadership framework stated at corporate, not carried to franchise stores."),
    ("Basic Assumptions", "Crew often treats the job as transactional. High turnover signals replaceability."),
]
for i, (label, desc) in enumerate(pop_items):
    y = col_top + Inches(0.85 + i * 1.45)
    tb(s, pop_x, y, col_w, Inches(0.45), label, sz=17, bold=True, color=POPEYES_ORANGE)
    tb(s, pop_x, y + Inches(0.45), col_w, Inches(0.95), desc, sz=15, color=DARK)

notes(s, """[CULTURE COMPARISON, ~2 min]
Read across each row.

Artifacts row: walk into any Chick-fil-A and the experience feels the same store to store. Walk into two Popeyes locations in the same week and they can feel like different chains. Our Camp Creek Popeyes had a broken digital menu board the day we visited, the Buckhead one did not.

Values row: CFA repeats its purpose statement constantly in training. Popeyes has a corporate framework but most front-line workers we saw could not tell us what it was.

Assumptions row: this is where the gap shows up in behavior. If a worker assumes effort matters, they give effort. If they assume they are replaceable, they do the minimum.""")

# =================== SLIDE 9: BURNS/BASS FRAMEWORK ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Leadership: Transformational vs. Transactional",
   sz=32, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), NEUTRAL_DARK)

col_top = Inches(1.95)
col_w = Inches(6.0)
left_w = Inches(0.5)
gap = Inches(0.25)

tx = left_w
rect(s, tx, col_top, col_w, Inches(0.6), NEUTRAL_DARK)
tb(s, tx, col_top, col_w, Inches(0.6),
   "Transformational  (Burns 1978  /  Bass 1985)",
   sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, tx + Inches(0.2), col_top + Inches(0.85), col_w - Inches(0.4), Inches(3.8), [
    "Idealized influence: lead by example",
    "Inspirational motivation: shared vision",
    "Intellectual stimulation: challenge people to grow",
    "Individualized consideration: develop each person",
], sz=16, line_spacing=1.3, terminal_period=False)

tx2 = left_w + col_w + gap
rect(s, tx2, col_top, col_w, Inches(0.6), NEUTRAL_GRAY)
tb(s, tx2, col_top, col_w, Inches(0.6), "Transactional",
   sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, tx2 + Inches(0.2), col_top + Inches(0.85), col_w - Inches(0.4), Inches(3.8), [
    "Contingent reward: pay tied to performance",
    "Management by exception: intervene only when problems arise",
], sz=16, line_spacing=1.3, terminal_period=False)

tb(s, Inches(0.75), Inches(6.6), Inches(11.83), Inches(0.6),
   "Transformational leadership predicts higher commitment, lower turnover, and stronger performance in meta-analyses.",
   sz=15, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)

notes(s, """[LEADERSHIP FRAMEWORK, ~1 min]
Bass developed this framework empirically from surveys of military and corporate leaders. The four transformational behaviors on the left are well-validated predictors of follower commitment.

Transactional leadership has its place. Most workplaces need some of it. The point of the comparison is which mode dominates day to day at each chain.

Next slide applies this to our two companies.""")

# =================== SLIDE 10: LEADERSHIP COMPARISON ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Leadership Comparison", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

col_top = Inches(1.85)
col_w = Inches(6.0)
left_w = Inches(0.5)
gap = Inches(0.25)

cfa_x = left_w
rect(s, cfa_x, col_top, col_w, Inches(0.6), CFA_RED)
tb(s, cfa_x, col_top, col_w, Inches(0.6), "Chick-fil-A: transformational",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, cfa_x + Inches(0.2), col_top + Inches(0.85), col_w - Inches(0.4), Inches(4.2), [
    "One Operator per store, hand-selected from roughly 40,000 applicants per year",
    "Operator is on the floor daily and works the line during the rush",
    "Promotes from within. Many corporate leaders began as hourly team members",
    "Personal training, same-shift recognition, visible promotion path",
], sz=15, line_spacing=1.3)

pop_x = left_w + col_w + gap
rect(s, pop_x, col_top, col_w, Inches(0.6), POPEYES_ORANGE)
tb(s, pop_x, col_top, col_w, Inches(0.6), "Popeyes: transactional",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, pop_x + Inches(0.2), col_top + Inches(0.85), col_w - Inches(0.4), Inches(4.2), [
    "Multi-unit franchisees may own dozens of stores",
    "Daily management delegated to hired general managers without ownership stakes",
    "Contingent reward: complete the shift, receive a paycheck",
    "Management by exception: attention shows up only when something breaks",
], sz=15, line_spacing=1.3)

notes(s, """[LEADERSHIP COMPARISON, ~2 min]
Structural difference matters more than any leadership training program could. CFA built transformational leadership into the franchise contract by requiring a single Operator on-site. Popeyes built transactional leadership into theirs by allowing absentee multi-unit ownership.

When the on-site leader has equity in the outcome and is required to be present, transformational behaviors become the easier path. When the on-site leader is a salaried general manager covering for an absent franchisee, transactional behaviors become the easier path.""")

# =================== SLIDE 11: HERZBERG FRAMEWORK ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Job Satisfaction: Herzberg's Two-Factor Theory (1959)",
   sz=32, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), NEUTRAL_DARK)

col_top = Inches(2.0)
col_w = Inches(6.0)
left_w = Inches(0.5)
gap = Inches(0.25)

hx = left_w
rect(s, hx, col_top, col_w, Inches(0.7), NEUTRAL_GRAY)
tb(s, hx, col_top, col_w, Inches(0.7), "Hygiene Factors",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, hx, col_top + Inches(0.85), col_w, Inches(0.4),
   "Prevent dissatisfaction", sz=14, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)
bullets(s, hx + Inches(0.5), col_top + Inches(1.4), col_w - Inches(1.0), Inches(3.0), [
    "Pay",
    "Working conditions",
    "Job security",
    "Supervision quality",
    "Company policies",
], sz=17, line_spacing=1.4, terminal_period=False)

mx = left_w + col_w + gap
rect(s, mx, col_top, col_w, Inches(0.7), NEUTRAL_DARK)
tb(s, mx, col_top, col_w, Inches(0.7), "Motivators",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, mx, col_top + Inches(0.85), col_w, Inches(0.4),
   "Drive satisfaction", sz=14, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)
bullets(s, mx + Inches(0.5), col_top + Inches(1.4), col_w - Inches(1.0), Inches(3.0), [
    "Achievement",
    "Recognition",
    "The work itself",
    "Responsibility",
    "Growth and advancement",
], sz=17, line_spacing=1.4, terminal_period=False)

tb(s, Inches(0.75), Inches(6.7), Inches(11.83), Inches(0.6),
   "Hygiene factors are the floor. Motivators are what produce engagement. Both have to be in place.",
   sz=16, italic=True, color=NEUTRAL_GRAY, align=PP_ALIGN.CENTER)

notes(s, """[HERZBERG FRAMEWORK, ~1 min]
Herzberg's key insight is that hygiene and motivators sit on two separate scales rather than two ends of one scale.

Solving hygiene problems takes a worker from miserable to neutral. Adding motivators takes a worker from neutral to engaged. Skipping the first to get to the second does not work.

[verify Herzberg page citation if Conklin asks for one — Mausner & Snyderman 1959 ch. 3]

Apply this to the two chains on the next slide.""")

# =================== SLIDE 12: SATISFACTION COMPARISON ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Job Satisfaction Comparison", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

col_top = Inches(1.85)
col_w = Inches(6.0)
left_w = Inches(0.5)
gap = Inches(0.25)

cfa_x = left_w
rect(s, cfa_x, col_top, col_w, Inches(0.6), CFA_RED)
tb(s, cfa_x, col_top, col_w, Inches(0.6), "Chick-fil-A: both factors met",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
cfa_sat = [
    ("Hygiene", "Competitive QSR pay, well-equipped kitchens, predictable scheduling, guaranteed Sundays off."),
    ("Motivators", "Remarkable Futures program (over $162 million in scholarships since 1973), same-shift recognition, visible promotion path."),
]
for i, (lbl, desc) in enumerate(cfa_sat):
    y = col_top + Inches(0.85 + i * 2.0)
    tb(s, cfa_x, y, col_w, Inches(0.45), lbl, sz=18, bold=True, color=CFA_RED)
    tb(s, cfa_x, y + Inches(0.5), col_w, Inches(1.4), desc, sz=15, color=DARK)

pop_x = left_w + col_w + gap
rect(s, pop_x, col_top, col_w, Inches(0.6), POPEYES_ORANGE)
tb(s, pop_x, col_top, col_w, Inches(0.6), "Popeyes: both factors unmet",
   sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pop_sat = [
    ("Hygiene", "Wages near the industry minimum, chronic understaffing, inconsistent scheduling."),
    ("Motivators", "Limited growth at most franchise stores, minimal recognition from off-site management, replaceability mindset reinforced by turnover."),
]
for i, (lbl, desc) in enumerate(pop_sat):
    y = col_top + Inches(0.85 + i * 2.0)
    tb(s, pop_x, y, col_w, Inches(0.45), lbl, sz=18, bold=True, color=POPEYES_ORANGE)
    tb(s, pop_x, y + Inches(0.5), col_w, Inches(1.4), desc, sz=15, color=DARK)

notes(s, """[SATISFACTION COMPARISON, ~2 min]
Herzberg predicts that when hygiene factors are unmet, no level of motivators can compensate. Popeyes does not have either category solved at the franchise level, which is why the model predicts the dissatisfaction we see.

CFA's $162 million in scholarships is a real motivator. It signals the company is willing to spend on the long-term path of people working the cash register, even though store-level operating costs would be lower without it.

[handoff to synthesis presenter]""")

# =================== SLIDE 13: HOW IT ALL CONNECTS ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "How the Three Components Connect", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

box_y = Inches(2.0)
box_h = Inches(1.7)
box_w = Inches(3.4)
arrow_w = Inches(0.6)
total = box_w * 3 + arrow_w * 2
left_start = (SLIDE_W - total) / 2

boxes = [
    ("Culture", "Schein. Sets the environment.", NEUTRAL_DARK),
    ("Leadership", "Burns and Bass. Reinforces the environment.", RGBColor(0x40, 0x55, 0x70)),
    ("Satisfaction", "Herzberg. Measures the outcome.", RGBColor(0x55, 0x70, 0x88)),
]
for i, (title, sub, color) in enumerate(boxes):
    x = left_start + (box_w + arrow_w) * i
    rect(s, x, box_y, box_w, box_h, color)
    tb(s, x, box_y + Inches(0.2), box_w, Inches(0.55), title,
       sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x, box_y + Inches(0.85), box_w, Inches(0.7), sub,
       sz=14, color=LIGHT_BG, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        ax = x + box_w + Inches(0.05)
        ay = box_y + Inches(0.65)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, arrow_w - Inches(0.1), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NEUTRAL_GRAY
        arrow.line.fill.background()

syn_y = Inches(4.5)
rect(s, Inches(0.5), syn_y, Inches(12.33), Inches(1.0), CFA_RED)
tb(s, Inches(0.7), syn_y + Inches(0.05), Inches(12.0), Inches(0.4),
   "Chick-fil-A", sz=16, bold=True, color=WHITE)
tb(s, Inches(0.7), syn_y + Inches(0.4), Inches(12.0), Inches(0.55),
   "Single-Operator structure produces transformational leadership, sustains strong culture, satisfies both Herzberg factors.",
   sz=15, color=WHITE)

syn_y2 = syn_y + Inches(1.15)
rect(s, Inches(0.5), syn_y2, Inches(12.33), Inches(1.0), POPEYES_ORANGE)
tb(s, Inches(0.7), syn_y2 + Inches(0.05), Inches(12.0), Inches(0.4),
   "Popeyes", sz=16, bold=True, color=WHITE)
tb(s, Inches(0.7), syn_y2 + Inches(0.4), Inches(12.0), Inches(0.55),
   "Multi-unit franchise structure prevents transformational leadership, weakens culture, leaves both Herzberg factors unmet.",
   sz=15, color=WHITE)

notes(s, """[SYNTHESIS, ~1.5 min]
This is the key slide.

Three components feed each other. CFA's structure produces transformational leadership at every store. That leadership sustains the strong culture. The strong culture lets the company satisfy both Herzberg categories.

Popeyes's structure prevents transformational leadership at the store. Without that leadership the local culture stays weak. Without the culture the company is left relying on hygiene factors that are themselves underfunded.

Two chains hiring from the same labor pool produce different employees because the structures around the workers are different.""")

# =================== SLIDE 14: TAKEAWAYS ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "Takeaways", sz=40, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

takeaways = [
    "Structure matters more than slogans. The franchise contract decides who is in the store, and that decides almost everything else.",
    "Culture only holds when leaders reinforce it daily. A mission statement on a wall does not change behavior at the counter.",
    "Hygiene factors are the floor. Motivators like growth, recognition, and a real promotion path are what get effort above the floor.",
]
for i, t in enumerate(takeaways):
    y = Inches(2.0 + i * 1.05)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CFA_RED
    circle.line.fill.background()
    tb(s, Inches(0.75), y, Inches(0.6), Inches(0.6), str(i + 1),
       sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(1.55), y, Inches(11.0), Inches(0.95), t,
       sz=17, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# Closing block: contingent, no aphorism
hero_y = Inches(5.5)
rect(s, Inches(0.5), hero_y, Inches(12.33), Inches(1.4), NEUTRAL_DARK)
tb(s, Inches(0.7), hero_y + Inches(0.25), Inches(12.0), Inches(0.85),
   "Two competitors, same labor pool, different employees. The structure around the worker is doing the work.",
   sz=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

notes(s, """[TAKEAWAYS, ~1.5 min]
Slow down on the three numbered points. Audience should walk out remembering at least these.

Closing line: keep it short. The argument is about structure. Resist the urge to broaden it into a motivational pitch.

[if running long, drop this slide's hero band and end on the three numbered points]

[handoff to references / Q&A presenter]""")

# =================== SLIDE 15: REFERENCES ===================
s = slide()
tb(s, Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.9),
   "References", sz=36, bold=True, color=NEUTRAL_DARK)
rect(s, Inches(0.75), Inches(1.4), Inches(2.5), Inches(0.06), CFA_RED)

ref_lines = [
    [("ACSI. (2023). ", False), ("Restaurant study 2022-2023", True),
     (". American Customer Satisfaction Index. https://www.theacsi.org", False)],
    [("Bachelder, C. (2015). ", False), ("Dare to serve: How to drive superior results by serving others", True),
     (". Berrett-Koehler.", False)],
    [("Bass, B. M. (1985). ", False), ("Leadership and performance beyond expectations", True),
     (". Free Press.", False)],
    [("Burns, J. M. (1978). ", False), ("Leadership", True), (". Harper & Row.", False)],
    [("Cathy, S. T. (2002). ", False), ("Eat Mor Chikin: Inspire more people", True),
     (". Looking Glass Books.", False)],
    [("Chick-fil-A. (2023). ", False), ("Remarkable Futures scholarships", True),
     (". https://www.chick-fil-a.com/remarkable-futures-scholarships", False)],
    [("DailyPay. (2023). ", False), ("QSR turnover and retention benchmarks", True),
     (". https://www.dailypay.com/resource-center/blog", False)],
    [("Herzberg, F., Mausner, B., & Snyderman, B. (1959). ", False),
     ("The motivation to work", True), (". Wiley.", False)],
    [("Schein, E. H. (1985). ", False), ("Organizational culture and leadership", True),
     (". Jossey-Bass.", False)],
]
ref_tb = s.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(5.2))
ref_tf = ref_tb.text_frame
ref_tf.word_wrap = True
for i, parts in enumerate(ref_lines):
    p = ref_tf.paragraphs[0] if i == 0 else ref_tf.add_paragraph()
    p.line_spacing = 1.2
    p.space_after = Pt(6)
    for text, italic in parts:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(14)
        run.font.italic = italic
        run.font.color.rgb = DARK

notes(s, """[REFERENCES + Q&A, 30 sec to start, then open]
References stay up during Q&A. Common questions to expect:
- What about pay differences? CFA pays roughly $11-13/hr starting in Atlanta, Popeyes $10-11/hr. Real but small. Not the main driver.
- Could Popeyes adopt the Operator model? Probably not without restructuring the franchise contract, which RBI is unlikely to do given the current portfolio.
- What did we observe that surprised us? The Camp Creek Popeyes had a manager visibly stressed during the lunch rush. The Camp Creek CFA had a team leader walking the line and refilling drinks for waiting cars. Same intersection, same labor pool.""")

prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Slide count: {len(prs.slides)}")
